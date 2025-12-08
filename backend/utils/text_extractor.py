import fitz
import easyocr
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import io
import os
from pptx import Presentation

# Legacy reader for backward compatibility
reader = easyocr.Reader(['en', 'hi'])

def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from PDF using PyMuPDF (fitz)
    LEGACY METHOD - Use OCRManager for new uploads with confidence scoring
    """
    text = ""
    doc = fitz.open(file_path)
    
    for page in doc:
        text += page.get_text()
        
        # If no text found, try OCR on images
        if not text.strip():
            pix = page.get_pixmap()
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            ocr_result = reader.readtext(img_data, detail=0)
            text += " ".join(ocr_result)
    
    doc.close()
    return text.strip()

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX"""
    doc = DocxDocument(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text.strip()

def extract_text_from_image(file_path: str) -> str:
    """Extract text from image using EasyOCR"""
    result = reader.readtext(file_path, detail=0)
    return " ".join(result)

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX"""
    prs = Presentation(file_path)
    text_parts = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
            
            # Extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            slide_text.append(cell.text)
        
        # Add slide separator
        if slide_text:
            text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
    
    return "\n\n".join(text_parts).strip()


def extract_text_from_txt(file_path: str) -> str:
    """Extract text from plain text file"""
    try:
        # Try UTF-8 first (most common)
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read().strip()
    except UnicodeDecodeError:
        # Fallback to latin-1 if UTF-8 fails
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read().strip()
        except Exception as e:
            raise ValueError(f"Failed to read text file: {str(e)}")

def extract_text(file_path: str, file_type: str) -> str:
    """
    Main extraction function
    LEGACY METHOD - Returns only text string
    Use extract_text_enhanced() for new uploads with OCR metadata
    """
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "docx":
        return extract_text_from_docx(file_path)
    elif file_type == "pptx":
        return extract_text_from_pptx(file_path)
    elif file_type == "txt":
        return extract_text_from_txt(file_path)
    elif file_type in ["jpeg", "jpg", "png"]:
        return extract_text_from_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def extract_text_enhanced(file_path: str, file_type: str, use_ocr=True):
    """
    Enhanced extraction with OCR support and metadata
    
    Args:
        file_path: Path to file
        file_type: File extension
        use_ocr: Whether to use enhanced OCR for scanned docs
        
    Returns:
        dict: {
            'text': str,
            'is_scanned': bool,
            'ocr_metadata': dict or None
        }
    """
    file_type = file_type.lower()
    
    # For non-OCR file types, use legacy extraction
    if file_type in ["docx", "pptx", "txt"]:
        if file_type == "docx":
            text = extract_text_from_docx(file_path)
        elif file_type == "pptx":
            text = extract_text_from_pptx(file_path)
        else:
            text = extract_text_from_txt(file_path)
            
        return {
            'text': text,
            'is_scanned': False,
            'ocr_metadata': None
        }
    
    # For PDF and images, use OCR Manager if enabled
    if use_ocr and file_type in ["pdf", "jpg", "jpeg", "png", "tiff", "bmp"]:
        try:
            from backend.utils.ocr import OCRManager
            
            ocr_manager = OCRManager(languages=['en', 'hi'])
            result = ocr_manager.extract_text(file_path, file_type, preprocessing_level='medium')
            
            return {
                'text': result['text'],
                'is_scanned': True,
                'ocr_metadata': result
            }
        except Exception as e:
            print(f"OCR extraction failed, falling back to legacy: {str(e)}")
            # Fallback to legacy extraction
            if file_type == "pdf":
                text = extract_text_from_pdf(file_path)
            else:
                text = extract_text_from_image(file_path)
                
            return {
                'text': text,
                'is_scanned': True,
                'ocr_metadata': None
            }
    else:
        # Use legacy extraction
        text = extract_text(file_path, file_type)
        return {
            'text': text,
            'is_scanned': False,
            'ocr_metadata': None
        }
